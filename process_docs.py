##!/usr/bin/env python
# coding: utf-8


from flask import Flask, request, jsonify
import os
import json
import re
import pandas as pd
import docx
from pptx import Presentation
import pdfplumber

app = Flask(__name__)

ALLOWED_EXTENSIONS = {'txt', 'pdf', 'csv', 'docx', 'xlsx', 'pptx'}

import re

#### Logger ####

# Crear el directorio de logs si no existe
if not os.path.exists('logs'):
    os.mkdir('logs')

# Configurar el manejador de logs para escribir en un archivo
log_file_path = 'logs/chatbotiurban.log'
file_handler = RotatingFileHandler(log_file_path, maxBytes=10240000, backupCount=1)
file_handler.setFormatter(logging.Formatter(
    '%(asctime)s %(levelname)s: %(message)s [in %(pathname)s:%(lineno)d]'
))
file_handler.setLevel(logging.DEBUG)  # Usa DEBUG o INFO según necesites

# Añadir el manejador de archivos al logger de la aplicación
app.logger.addHandler(file_handler)

# También añadir un manejador de consola para la salida estándar
console_handler = logging.StreamHandler()
console_handler.setFormatter(logging.Formatter(
    '%(asctime)s %(levelname)s: %(message)s [in %(pathname)s:%(lineno)d]'
))
console_handler.setLevel(logging.DEBUG)  # Asegúrate de que este nivel sea consistente con file_handler.setLevel
app.logger.addHandler(console_handler)

# Establecer el nivel del logger de la aplicación
app.logger.setLevel(logging.DEBUG)

app.logger.info('Inicio de la aplicación ChatbotIUrban')

#### FIN Logger ####

# Configuración del Logging
def setup_logging():
    if not os.path.exists('logs'):
        os.mkdir('logs')

    log_file_path = 'logs/chatbotiurban.log'
    file_handler = RotatingFileHandler(log_file_path, maxBytes=10240000, backupCount=1)
    file_handler.setFormatter(logging.Formatter(
        '%(asctime)s %(levelname)s: %(message)s [in %(pathname)s:%(lineno)d]'
    ))
    file_handler.setLevel(logging.DEBUG)

    console_handler = logging.StreamHandler()
    console_handler.setFormatter(logging.Formatter(
        '%(asctime)s %(levelname)s: %(message)s [in %(pathname)s:%(lineno)d]'
    ))
    console_handler.setLevel(logging.DEBUG)

    logger = logging.getLogger(__name__)
    logger.addHandler(file_handler)
    logger.addHandler(console_handler)
    logger.setLevel(logging.DEBUG)
    return logger

logger = setup_logging()

# Función de limpieza y formato de contenido
def clean_and_format_content(content):
    # ... misma función que proporcionaste anteriormente ...

# Función para leer archivos TXT con múltiples codificaciones
def read_txt(file_path):
    encodings = ['utf-8', 'latin-1', 'ISO-8859-1', 'windows-1252']
    for encoding in encodings:
        try:
            with open(file_path, 'r', encoding=encoding) as file:
                return clean_and_format_content(file.read())
        except UnicodeDecodeError:
            continue
    try:
        with open(file_path, 'rb') as file:
            return clean_and_format_content(file.read().decode('utf-8', 'ignore'))
    except Exception as e:
        logger.error(f"Error al leer archivo TXT: {e}")
        return f"Error al procesar archivo TXT: {e}"

# Función para leer archivos PDF
def read_pdf(file_path):
    text = ""
    try:
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text
        return clean_and_format_content(text.encode().decode('utf-8', 'ignore'))
    except Exception as e:
        logger.error(f"Error al procesar PDF: {e}")
        return f"Error al procesar PDF: {e}"

# Función para leer archivos CSV o XLSX
def read_csv_or_xlsx(file_path, extension):
    try:
        if extension == 'csv':
            df = pd.read_csv(file_path, encoding='utf-8', error_bad_lines=False)
        else:
            df = pd.read_excel(file_path)
        return clean_and_format_content(df.to_string())
    except Exception as e:
        logger.error(f"Error al procesar {extension.upper()} archivo: {e}")
        return f"Error al procesar archivo {extension.upper()}: {e}"

# Función para leer archivos DOCX
def read_docx(file_path):
    try:
        doc = docx.Document(file_path)
        content = "\n".join([para.text for para in doc.paragraphs])
        return clean_and_format_content(content)
    except Exception as e:
        logger.error(f"Error al procesar DOCX: {e}")
        return f"Error al procesar DOCX: {e}"

# Función para leer archivos PPTX
def read_pptx(file_path):
    try:
        ppt = Presentation(file_path)
        text = ""
        for slide in ppt.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return clean_and_format_content(text)
    except Exception as e:
        logger.error(f"Error al procesar PPTX: {e}")
        return f"Error al procesar PPTX: {e}"

# Función para procesar el archivo según su extensión
def process_file(file_path, extension):
    if extension == 'txt':
        return read_txt(file_path)
    elif extension == 'pdf':
        return read_pdf(file_path)
    elif extension in ['csv', 'xlsx']:
        return read_csv_or_xlsx(file_path, extension)
    elif extension == 'docx':
        return read_docx(file_path)
    elif extension == 'pptx':
        return read_pptx(file_path)
    else:
        logger.error(f"Formato de archivo no soportado: {extension}")
        return "Formato de archivo no soportado"
